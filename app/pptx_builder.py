"""Render a Deck into an editable .pptx file using python-pptx.

Bulleted slides use the standard title+content layout. Flow diagrams are drawn
as native, editable rounded-rectangle shapes connected by arrow connectors, so
the user can move/restyle them in PowerPoint afterwards.
"""

from __future__ import annotations

import io
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from .models import Deck, Diagram, Slide, Table
from .themes import DEFAULT_THEME, Theme


def _rgb(hex_str: str) -> RGBColor:
    """Convert a ``RRGGBB`` hex string into a python-pptx ``RGBColor``."""
    return RGBColor.from_string(hex_str)


def _slide_size(prs: Presentation) -> tuple[int, int]:
    return prs.slide_width, prs.slide_height


def _add_page_background(slide, prs: Presentation, theme: Theme) -> None:
    """Fill the whole slide with the theme background (needed for dark themes)."""
    if theme.page_bg.upper() in ("FFFFFF", "FFF"):
        return  # default white slide; nothing to paint
    width, height = _slide_size(prs)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, width, height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(theme.page_bg)
    bg.line.fill.background()
    bg.shadow.inherit = False


def _add_title_slide(prs: Presentation, deck: Deck, theme: Theme) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    width, height = _slide_size(prs)
    _add_page_background(slide, prs, theme)

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, int(height * 0.34), width, int(height * 0.32)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = _rgb(theme.accent_dark)
    band.line.fill.background()

    box = slide.shapes.add_textbox(
        int(width * 0.08), int(height * 0.36), int(width * 0.84), int(height * 0.28)
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = deck.title
    p.alignment = PP_ALIGN.LEFT
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = _rgb(theme.title_text)
    if deck.subtitle:
        sub = tf.add_paragraph()
        sub.text = deck.subtitle
        sub.font.size = Pt(18)
        sub.font.color.rgb = _rgb(theme.subtitle_text)


def _add_bullets(slide, slide_data: Slide, prs: Presentation, theme: Theme) -> None:
    width, _ = _slide_size(prs)
    box = slide.shapes.add_textbox(
        int(width * 0.08), Emu(int(prs.slide_height * 0.26)),
        int(width * 0.84), int(prs.slide_height * 0.6),
    )
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(slide_data.bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = _rgb(theme.body_text)
        p.space_after = Pt(10)


def _add_diagram(
    slide,
    diagram: Diagram,
    prs: Presentation,
    theme: Theme,
    roomy: bool = False,
) -> None:
    width, height = _slide_size(prs)
    steps = diagram.steps
    if not steps:
        return

    top_margin = int(height * (0.20 if roomy else 0.28))
    area_w = int(width * (0.90 if roomy else 0.84))
    left_margin = int(width * (0.05 if roomy else 0.08))

    if diagram.direction == "right":
        gap = int(area_w * (0.03 if roomy else 0.04))
        node_w = (area_w - gap * (len(steps) - 1)) // len(steps)
        node_h = int(height * (0.20 if roomy else 0.16))
        top = top_margin + int(height * (0.08 if roomy else 0.12))
        boxes = []
        for i, label in enumerate(steps):
            left = left_margin + i * (node_w + gap)
            boxes.append(_node(slide, label, left, top, node_w, node_h, theme))
        for a, b in zip(boxes, boxes[1:]):
            _connect(slide, a, b, theme)
    else:  # "down"
        node_w = int(area_w * (0.68 if roomy else 0.5))
        left = left_margin + (area_w - node_w) // 2
        available = int(height * (0.72 if roomy else 0.62))
        gap = int(available * (0.04 if roomy else 0.05))
        node_h = (available - gap * (len(steps) - 1)) // len(steps)
        node_h = min(node_h, int(height * (0.16 if roomy else 0.13)))
        boxes = []
        for i, label in enumerate(steps):
            top = top_margin + i * (node_h + gap)
            boxes.append(_node(slide, label, left, top, node_w, node_h, theme))
        for a, b in zip(boxes, boxes[1:]):
            _connect(slide, a, b, theme)


def _node(slide, label: str, left: int, top: int, w: int, h: int, theme: Theme):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(theme.node_fill)
    shape.line.color.rgb = _rgb(theme.accent)
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = _rgb(theme.node_text)
    return shape


def _connect(slide, a, b, theme: Theme) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        a.left + a.width // 2, a.top + a.height,
        b.left + b.width // 2, b.top,
    )
    # Vertical vs horizontal endpoints depending on relative position.
    if abs(b.top - a.top) < a.height:  # side by side -> horizontal arrow
        connector.begin_x, connector.begin_y = a.left + a.width, a.top + a.height // 2
        connector.end_x, connector.end_y = b.left, b.top + b.height // 2
    line = connector.line
    line.color.rgb = _rgb(theme.accent)
    line.width = Pt(2.25)
    _add_arrow_head(connector)


def _add_arrow_head(connector) -> None:
    """Add a triangle arrowhead to a connector's end (python-pptx has no helper)."""
    ln = connector.line._get_or_add_ln()
    from pptx.oxml.ns import qn

    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = ln.makeelement(qn("a:tailEnd"), {})
        ln.append(tail)
    tail.set("type", "triangle")
    tail.set("w", "med")
    tail.set("len", "med")


def _add_table(slide, table: Table, prs: Presentation, theme: Theme) -> None:
    width, height = _slide_size(prs)
    headers = table.headers or [""] * (len(table.rows[0]) if table.rows else 0)
    n_cols = max(len(headers), max((len(r) for r in table.rows), default=0))
    n_rows = len(table.rows) + 1  # + header row
    if n_cols == 0 or n_rows <= 1:
        return

    left = int(width * 0.06)
    top = int(height * 0.24)
    table_w = int(width * 0.88)
    table_h = int(height * 0.66)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, table_w, table_h)
    tbl = shape.table

    # Font shrinks as the table gets denser so more content stays on one slide.
    longest = max(
        (len(c) for row in table.rows for c in row),
        default=10,
    )
    font_pt = 12 if longest <= 40 else 10 if longest <= 90 else 8

    for col_idx in range(n_cols):
        header_text = headers[col_idx] if col_idx < len(headers) else ""
        cell = tbl.cell(0, col_idx)
        _style_cell(cell, header_text, font_pt, theme, header=True)

    for r_idx, row in enumerate(table.rows, start=1):
        for col_idx in range(n_cols):
            text = row[col_idx] if col_idx < len(row) else ""
            _style_cell(tbl.cell(r_idx, col_idx), text, font_pt, theme, header=False)


def _style_cell(cell, text: str, font_pt: int, theme: Theme, header: bool) -> None:
    cell.margin_left = Emu(45720)
    cell.margin_right = Emu(45720)
    cell.margin_top = Emu(18288)
    cell.margin_bottom = Emu(18288)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(theme.accent_dark if header else theme.table_cell_bg)
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_pt)
    p.font.bold = header
    p.font.color.rgb = _rgb(theme.table_header_text if header else theme.body_text)


def _add_content_slide(prs: Presentation, slide_data: Slide, theme: Theme) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    width, _ = _slide_size(prs)
    _add_page_background(slide, prs, theme)

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, int(width * 0.025), prs.slide_height
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(theme.accent)
    accent.line.fill.background()

    title_box = slide.shapes.add_textbox(
        int(width * 0.08), int(prs.slide_height * 0.08),
        int(width * 0.84), int(prs.slide_height * 0.14),
    )
    tp = title_box.text_frame.paragraphs[0]
    tp.text = slide_data.title
    tp.font.size = Pt(30)
    tp.font.bold = True
    tp.font.color.rgb = _rgb(theme.heading_text)

    is_diagram_only = bool(slide_data.diagram and not slide_data.bullets and not slide_data.table)
    if slide_data.bullets:
        _add_bullets(slide, slide_data, prs, theme)
    if slide_data.diagram:
        _add_diagram(slide, slide_data.diagram, prs, theme, roomy=is_diagram_only)
    if slide_data.table:
        _add_table(slide, slide_data.table, prs, theme)

    if slide_data.notes:
        slide.notes_slide.notes_text_frame.text = slide_data.notes


def _build_presentation(deck: Deck, theme: Theme) -> Presentation:
    prs = Presentation()
    prs.slide_width = Emu(12192000)  # 16:9
    prs.slide_height = Emu(6858000)
    _add_title_slide(prs, deck, theme)
    for slide_data in deck.slides:
        _add_content_slide(prs, slide_data, theme)
    return prs


def render_pptx(deck: Deck, theme: Theme = DEFAULT_THEME) -> bytes:
    """Build the deck entirely in memory (stateless / serverless-friendly)."""
    buffer = io.BytesIO()
    _build_presentation(deck, theme).save(buffer)
    return buffer.getvalue()


def build_pptx(deck: Deck, out_path: Path, theme: Theme = DEFAULT_THEME) -> Path:
    """Convenience for local use / CLI: render and write to disk."""
    out_path.write_bytes(render_pptx(deck, theme))
    return out_path
